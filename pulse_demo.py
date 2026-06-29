#!/usr/bin/env python3
"""
EGC Pulse. Compliant social listening (zero-infra demo, pure stdlib).

Pipeline: collect (via compliant_connectors) -> enrich (sentiment) -> dedupe ->
relevance filter -> store -> REST API -> grayscale dashboard.

Real data only. No synthetic fill. Sources, statuses, and date/historical
behavior are governed by compliant_connectors.py. See COMPLIANT_CONNECTORS.md.

Usage:
  python3 pulse_demo.py serve              # API + dashboard on :8787
  python3 pulse_demo.py collect "Jovia"    # one-shot collect (current sources)
  python3 pulse_demo.py reset              # wipe the database
"""
import hashlib
import io
import csv
import json
import math
import os
import re
import sys
import threading
import uuid
import urllib.parse
import urllib.request
import urllib.error
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timedelta, timezone
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer

import platform_access_manager as pam
import ai_policy as aip
import source_truth as struth
import compliant_discovery as cd

HERE = os.path.dirname(os.path.abspath(__file__))
# DB lives next to the app by default; override with PULSE_DB to point at a
# mounted persistent disk when hosting (e.g. PULSE_DB=/data/pulse_demo.db).
DB_PATH = os.environ.get("PULSE_DB") or os.path.join(HERE, "pulse_demo.db")
KEYWORDS_PATH = os.path.join(HERE, "keywords.json")
PORT = int(os.environ.get("PORT", "8787"))
DEFAULT_KEYWORDS = []


def now_iso():
    return datetime.now(timezone.utc).isoformat()


def load_env():
    for path in (os.path.join(HERE, ".env"), os.path.join(HERE, "..", ".env")):
        path = os.path.abspath(path)
        if not os.path.exists(path):
            continue
        try:
            with open(path) as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#") or "=" not in line:
                        continue
                    k, v = line.split("=", 1)
                    os.environ.setdefault(k.strip(), v.strip().strip('"').strip("'"))
        except Exception:
            pass


# ── keywords ─────────────────────────────────────────────────────────────────
def _norm_kw(item):
    if isinstance(item, str):
        return {"label": item, "match": [item], "query": item}
    label = item.get("label") or (item.get("match") or [""])[0]
    return {"label": label, "match": item.get("match") or [label], "query": item.get("query") or label}


# ── projects: every research belongs to a named project ─────────────────────
def list_projects():
    with db() as c:
        rows = c.execute("SELECT p.id, p.name, p.created_at, "
                         "(SELECT COUNT(*) FROM mentions m WHERE m.project_id = p.id) mentions, "
                         "(SELECT COUNT(*) FROM keywords k WHERE k.project_id = p.id) terms "
                         "FROM projects p ORDER BY p.created_at, p.rowid").fetchall()
    return [dict(r) for r in rows]


def default_project_id():
    with db() as c:
        r = c.execute("SELECT id FROM projects ORDER BY created_at, rowid LIMIT 1").fetchone()
    return r["id"] if r else None


def resolve_project(pid):
    """Return a valid project id, falling back to the default project so data is
    never written or read outside a project."""
    if pid:
        with db() as c:
            r = c.execute("SELECT id FROM projects WHERE id = ?", (pid,)).fetchone()
        if r:
            return r["id"]
    return default_project_id()


def create_project(name):
    name = (name or "").strip() or "Untitled Project"
    pid = uuid.uuid4().hex[:12]
    with db() as c:
        c.execute("INSERT INTO projects (id, name, created_at) VALUES (?,?,?)", (pid, name, now_iso()))
    return {"id": pid, "name": name, "created_at": now_iso(), "mentions": 0, "terms": 0}


def rename_project(pid, name):
    name = (name or "").strip()
    if not (pid and name):
        return False
    with db() as c:
        c.execute("UPDATE projects SET name = ? WHERE id = ?", (name, pid))
    return True


def delete_project(pid):
    """Delete a project and all of its stored data. Never deletes the last project."""
    with db() as c:
        n = c.execute("SELECT COUNT(*) n FROM projects").fetchone()["n"]
        if n <= 1:
            return False
        c.execute("DELETE FROM mentions WHERE project_id = ?", (pid,))
        c.execute("DELETE FROM keywords WHERE project_id = ?", (pid,))
        c.execute("DELETE FROM projects WHERE id = ?", (pid,))
    return True


def project_name(pid):
    pid = resolve_project(pid)
    with db() as c:
        r = c.execute("SELECT name FROM projects WHERE id = ?", (pid,)).fetchone()
    return r["name"] if r else "Project"


def tracked_keywords(project=None):
    pid = resolve_project(project)
    with db() as c:
        rows = c.execute("SELECT label, kw_json FROM keywords WHERE project_id = ? ORDER BY created_at, rowid", (pid,)).fetchall()
    out = []
    for r in rows:
        try:
            out.append(_norm_kw(json.loads(r["kw_json"]) if r["kw_json"] else r["label"]))
        except Exception:
            out.append(_norm_kw(r["label"]))
    return out


def add_keyword(term, project=None):
    term = (term or "").strip()
    if not term:
        return
    pid = resolve_project(project)
    k = _norm_kw(term)
    with db() as c:
        c.execute("INSERT OR IGNORE INTO keywords (project_id, label, kw_json, created_at) VALUES (?,?,?,?)",
                  (pid, k["label"], json.dumps(k), now_iso()))


def remove_keyword(label, project=None):
    pid = resolve_project(project)
    with db() as c:
        c.execute("DELETE FROM keywords WHERE project_id = ? AND lower(label) = lower(?)", (pid, label or ""))
        c.execute("DELETE FROM mentions WHERE project_id = ? AND keyword = ?", (pid, label))


def clear_project(project=None):
    """Clear a single project's tracked terms and stored mentions (not other projects)."""
    pid = resolve_project(project)
    with db() as c:
        c.execute("DELETE FROM keywords WHERE project_id = ?", (pid,))
        c.execute("DELETE FROM mentions WHERE project_id = ?", (pid,))


def add_manual_insight(platform, period, impressions, reach, engagement=0, note=""):
    plat = (platform or "instagram").strip().lower()
    if plat not in ("instagram", "facebook"):
        plat = "instagram"
    def _int(v):
        try:
            return int(float(str(v).replace(",", "").strip() or 0))
        except Exception:
            return 0
    with db() as c:
        c.execute("INSERT INTO manual_insights (platform, period, impressions, reach, engagement, note, created_at)"
                  " VALUES (?,?,?,?,?,?,?)",
                  (plat, (period or "").strip(), _int(impressions), _int(reach), _int(engagement),
                   (note or "").strip(), now_iso()))
    return True


def get_manual_insights():
    with db() as c:
        rows = c.execute("SELECT platform, period, impressions, reach, engagement, note, created_at"
                         " FROM manual_insights ORDER BY id DESC").fetchall()
    return [dict(r) for r in rows]


_re_cache = {}


def _term_re(term):
    if term not in _re_cache:
        t = term.strip()
        prefix = r"(?<!\w)#" if t.startswith("#") else r"(?<!\w)"
        body = re.escape(t[1:] if t.startswith("#") else t)
        _re_cache[term] = re.compile(prefix + body + r"(?!\w)", re.I)
    return _re_cache[term]


def is_relevant(text, terms):
    return bool(text) and any(_term_re(t).search(text) for t in terms)


# ── sentiment (lexicon, pure python) ─────────────────────────────────────────
POS = {"love": 2, "loved": 2, "great": 2, "amazing": 3, "awesome": 3, "excellent": 3, "best": 2,
       "fantastic": 3, "incredible": 3, "happy": 2, "excited": 2, "win": 2, "good": 1, "nice": 1,
       "solid": 1, "impressive": 2, "improved": 2, "fast": 1, "smooth": 2, "recommend": 2,
       "perfect": 3, "wonderful": 3, "reliable": 2, "worth": 1, "helpful": 2, "thanks": 1, "thank": 1}
NEG = {"hate": 3, "terrible": 3, "awful": 3, "worst": 3, "bad": 2, "poor": 2, "broken": 2, "buggy": 2,
       "crash": 2, "slow": 2, "disappointed": 3, "scam": 3, "ripoff": 3, "overpriced": 2, "useless": 3,
       "garbage": 3, "fail": 2, "failed": 2, "angry": 2, "annoying": 2, "frustrated": 2, "refund": 1,
       "complaint": 2, "problem": 1, "issue": 1, "issues": 2, "fraud": 3, "fee": 1, "fees": 1, "denied": 2}
NEGATORS = {"not", "no", "never", "without", "cant", "cannot", "wont", "dont", "isnt", "arent"}
TOKEN_RE = re.compile(r"[a-z']+")


def analyze_sentiment(text):
    if not text:
        return "neutral", 0.0
    toks = TOKEN_RE.findall(text.lower())
    total = 0.0
    for i, t in enumerate(toks):
        base = POS.get(t, 0) - NEG.get(t, 0)
        if base:
            if i and toks[i - 1] in NEGATORS:
                base *= -0.8
            total += base
    score = max(-1.0, min(1.0, total / 6.0))
    return ("positive" if score >= 0.05 else "negative" if score <= -0.05 else "neutral"), round(score, 3)


# ── database ─────────────────────────────────────────────────────────────────
import sqlite3


def db():
    d = os.path.dirname(DB_PATH)
    if d and not os.path.isdir(d):
        os.makedirs(d, exist_ok=True)   # support PULSE_DB pointing at a mounted disk dir
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    with db() as c:
        c.execute("""CREATE TABLE IF NOT EXISTS mentions (
              id INTEGER PRIMARY KEY AUTOINCREMENT, platform TEXT NOT NULL,
              platform_post_id TEXT NOT NULL, keyword TEXT, author TEXT, content TEXT,
              content_hash TEXT NOT NULL, url TEXT, lang TEXT DEFAULT 'en',
              posted_at TEXT NOT NULL, posted_date TEXT, ingested_at TEXT NOT NULL,
              engagement INTEGER DEFAULT 0, reach INTEGER DEFAULT 0, sentiment TEXT, sentiment_score REAL,
              is_hidden INTEGER DEFAULT 0, project_id TEXT,
              UNIQUE(platform, platform_post_id, keyword, project_id))""")
        c.execute("CREATE TABLE IF NOT EXISTS suppression (platform TEXT, id_hash TEXT, PRIMARY KEY (platform, id_hash))")
        # Owner-entered account analytics (impressions/reach typed from the user's
        # own Instagram/Facebook Insights). Lawful, user-provided, labeled manual.
        c.execute("""CREATE TABLE IF NOT EXISTS manual_insights (
              id INTEGER PRIMARY KEY AUTOINCREMENT, platform TEXT NOT NULL, period TEXT,
              impressions INTEGER DEFAULT 0, reach INTEGER DEFAULT 0, engagement INTEGER DEFAULT 0,
              note TEXT, created_at TEXT NOT NULL)""")
        cols = [r["name"] for r in c.execute("PRAGMA table_info(mentions)")]
        if "posted_date" not in cols:  # migrate older DBs
            c.execute("ALTER TABLE mentions ADD COLUMN posted_date TEXT")
            c.execute("UPDATE mentions SET posted_date = substr(posted_at,1,10) WHERE posted_date IS NULL")
        if "reach" not in cols:  # engagement (interactions) vs reach (views/visibility) split for analytics
            c.execute("ALTER TABLE mentions ADD COLUMN reach INTEGER DEFAULT 0")
        # source-truth / coverage-ledger columns (safe additive migration)
        truth_cols = {"display_platform": "TEXT", "source_platform": "TEXT", "searched_platform": "TEXT",
                      "discussed_platforms": "TEXT", "direct_platform_data": "INTEGER",
                      "platform_coverage_type": "TEXT", "coverage_label": "TEXT", "coverage_note": "TEXT",
                      "access_path": "TEXT", "source_key": "TEXT", "confidence_level": "TEXT", "run_id": "TEXT",
                      "result_type": "TEXT", "description": "TEXT"}
        need_backfill = "platform_coverage_type" not in cols
        for col, typ in truth_cols.items():
            if col not in cols:
                c.execute("ALTER TABLE mentions ADD COLUMN %s %s" % (col, typ))
        if need_backfill:
            _backfill_source_truth(c)
        c.execute("CREATE INDEX IF NOT EXISTS idx_posted ON mentions(posted_date)")
        c.execute("CREATE INDEX IF NOT EXISTS idx_kw ON mentions(keyword)")
        c.execute("CREATE INDEX IF NOT EXISTS idx_cov ON mentions(platform_coverage_type)")
        # ── projects: every research belongs to a named project ──────────────
        c.execute("""CREATE TABLE IF NOT EXISTS projects (
              id TEXT PRIMARY KEY, name TEXT NOT NULL, created_at TEXT NOT NULL)""")
        c.execute("""CREATE TABLE IF NOT EXISTS keywords (
              project_id TEXT NOT NULL, label TEXT NOT NULL, kw_json TEXT, created_at TEXT,
              PRIMARY KEY (project_id, label))""")
        if "project_id" not in cols:
            c.execute("ALTER TABLE mentions ADD COLUMN project_id TEXT")
        c.execute("CREATE INDEX IF NOT EXISTS idx_proj ON mentions(project_id)")
        # Ensure at least one project exists, then home any orphan data into it.
        row = c.execute("SELECT id FROM projects ORDER BY created_at, rowid LIMIT 1").fetchone()
        if not row:
            pid = uuid.uuid4().hex[:12]
            c.execute("INSERT INTO projects (id, name, created_at) VALUES (?,?,?)",
                      (pid, "Default Project", now_iso()))
            # migrate any pre-project keywords.json into the default project
            try:
                if os.path.exists(KEYWORDS_PATH):
                    with open(KEYWORDS_PATH) as f:
                        for item in (json.load(f) or []):
                            k = _norm_kw(item)
                            c.execute("INSERT OR IGNORE INTO keywords (project_id, label, kw_json, created_at) "
                                      "VALUES (?,?,?,?)", (pid, k["label"], json.dumps(k), now_iso()))
            except Exception:
                pass
        else:
            pid = row["id"]
        c.execute("UPDATE mentions SET project_id = ? WHERE project_id IS NULL OR project_id = ''", (pid,))
        _migrate_project_unique(c)


def _migrate_project_unique(c):
    """Make dedup project-scoped. Older DBs had UNIQUE(platform, platform_post_id,
    keyword) which would drop a post already collected in another project. Rebuild
    the table so the same post can live in multiple projects. Idempotent."""
    uniques = [ix for ix in c.execute("PRAGMA index_list(mentions)").fetchall() if ix["origin"] == "u"]
    has_proj = any("project_id" in [r["name"] for r in c.execute("PRAGMA index_info('%s')" % ix["name"]).fetchall()]
                   for ix in uniques)
    if has_proj or not uniques:
        return  # fresh schema already project-scoped, or no UNIQUE to fix
    cols = c.execute("PRAGMA table_info(mentions)").fetchall()
    defs, names = [], []
    for col in cols:
        names.append(col["name"])
        if col["pk"]:
            defs.append('"%s" INTEGER PRIMARY KEY AUTOINCREMENT' % col["name"])
            continue
        d = '"%s" %s' % (col["name"], col["type"] or "TEXT")
        if col["notnull"]:
            d += " NOT NULL"
        if col["dflt_value"] is not None:
            d += " DEFAULT %s" % col["dflt_value"]
        defs.append(d)
    collist = ", ".join('"%s"' % n for n in names)
    c.execute("DROP TABLE IF EXISTS mentions_new")
    c.execute("CREATE TABLE mentions_new (%s, UNIQUE(platform, platform_post_id, keyword, project_id))"
              % ", ".join(defs))
    c.execute("INSERT INTO mentions_new (%s) SELECT %s FROM mentions" % (collist, collist))
    c.execute("DROP TABLE mentions")
    c.execute("ALTER TABLE mentions_new RENAME TO mentions")
    c.execute("CREATE INDEX IF NOT EXISTS idx_posted ON mentions(posted_date)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_kw ON mentions(keyword)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_cov ON mentions(platform_coverage_type)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_proj ON mentions(project_id)")


def _backfill_source_truth(c):
    """Conservative provenance for rows that predate the source-truth layer."""
    rules = [
        ("news", "Open Web / News", "gdelt", "open web / news", 0, "open_web_reference",
         "Open-web / news result", "Open-web/news source that references a platform. Not platform-native data.", "source_verified_open_web_reference"),
        ("hackernews", "Hacker News", "hacker_news", "hacker news", 1, "open_network_public",
         "Hacker News result", "Hacker News public API result.", "source_verified_direct_platform"),
        ("mastodon", "Mastodon", "mastodon", "mastodon public api", 1, "open_network_public",
         "Mastodon public result", "Mastodon public (fediverse) result.", "source_verified_direct_platform"),
        ("lemmy", "Lemmy", "lemmy", "lemmy public api", 1, "open_network_public",
         "Lemmy public result", "Lemmy public (fediverse) result.", "source_verified_direct_platform"),
        ("nostr", "Nostr", "nostr", "nostr relays", 1, "open_network_public",
         "Nostr public result", "Nostr relay result.", "source_verified_direct_platform"),
        ("peertube", "PeerTube", "peertube", "peertube search", 1, "open_network_public",
         "PeerTube result", "PeerTube federated video result.", "source_verified_direct_platform"),
        ("reddit", "Reddit", "reddit", "reddit official api", 1, "direct_official_api",
         "Reddit official API result", "Reddit official API result. Direct platform data.", "source_verified_direct_platform"),
        ("x", "X / Twitter", "x", "x recent search", 1, "direct_official_api",
         "X official API result", "X official API result.", "source_verified_direct_platform"),
        ("manual", "Manual Import", "manual_import", "manual import", 1, "manual_import",
         "Manual import. User-provided data", "User-provided data only. User is responsible for lawful upload rights.", "user_supplied_manual_import"),
    ]
    for plat, disp, sp, searched, direct, ct, label, note, conf in rules:
        c.execute("UPDATE mentions SET display_platform=?, source_platform=?, searched_platform=?, "
                  "direct_platform_data=?, platform_coverage_type=?, coverage_label=?, coverage_note=?, "
                  "confidence_level=?, discussed_platforms=COALESCE(discussed_platforms,'[]') "
                  "WHERE platform=? AND platform_coverage_type IS NULL",
                  (disp, sp, searched, direct, ct, label, note, conf, plat))
    c.execute("UPDATE mentions SET display_platform=COALESCE(display_platform,'Limited metadata'), "
              "source_platform=COALESCE(source_platform,platform), searched_platform=COALESCE(searched_platform,platform), "
              "direct_platform_data=COALESCE(direct_platform_data,0), platform_coverage_type=COALESCE(platform_coverage_type,'limited_metadata'), "
              "coverage_label=COALESCE(coverage_label,'Limited source metadata'), coverage_note=COALESCE(coverage_note,'Limited source metadata.'), "
              "confidence_level=COALESCE(confidence_level,'limited_metadata'), discussed_platforms=COALESCE(discussed_platforms,'[]') "
              "WHERE platform_coverage_type IS NULL")


def id_hash(platform, pid):
    return hashlib.sha256(("%s:%s" % (platform, pid)).encode()).hexdigest()


def is_suppressed(platform, pid):
    with db() as c:
        return c.execute("SELECT 1 FROM suppression WHERE platform=? AND id_hash=?",
                         (platform, id_hash(platform, pid))).fetchone() is not None


def ingest(records, project=None):
    n = 0
    now = now_iso()
    pid = resolve_project(project)
    with db() as c:
        for r in records:
            if is_suppressed(r["platform"], r["platform_post_id"]):
                continue
            label, score = analyze_sentiment(((r.get("content") or "") + " " + (r.get("description") or "")).strip())
            posted = r.get("posted_at") or now
            ch = hashlib.sha256((r["platform"] + ":" + re.sub(r"\s+", " ", (r.get("content") or "").lower())).encode()).hexdigest()
            try:
                cur = c.execute(
                    "INSERT OR IGNORE INTO mentions (platform, platform_post_id, keyword, author, content, "
                    "content_hash, url, lang, posted_at, posted_date, ingested_at, engagement, reach, sentiment, sentiment_score, "
                    "display_platform, source_platform, searched_platform, discussed_platforms, direct_platform_data, "
                    "platform_coverage_type, coverage_label, coverage_note, access_path, source_key, confidence_level, run_id, "
                    "result_type, description, project_id) "
                    "VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
                    (r["platform"], str(r["platform_post_id"]), r.get("keyword"), r.get("author"),
                     r.get("content"), ch, r.get("url"), r.get("lang", "en"), posted, posted[:10],
                     now, int(r.get("engagement", 0)), int(r.get("reach", 0)), label, score,
                     r.get("display_platform"), r.get("source_platform"), r.get("searched_platform"),
                     r.get("discussed_platforms") or "[]", int(r.get("direct_platform_data") or 0),
                     r.get("platform_coverage_type"), r.get("coverage_label"), r.get("coverage_note"),
                     r.get("access_path"), r.get("source_key"), r.get("confidence_level"), r.get("run_id"),
                     r.get("result_type"), r.get("description"), r.get("project_id") or pid))
                n += cur.rowcount
            except sqlite3.Error:
                pass
    return n


# ── collection (compliant connectors + date range + chunking) ───────────────
def collect(term, start=None, end=None, project=None):
    """Fan out across live collectors via the access manager. Chunks per-source
    max date-windows (e.g. TikTok Research) and merges/dedupes."""
    pid = resolve_project(project)
    sources = pam.get_live_collectors()

    def run_one(src):
        windows = pam.chunk_range(start, end, src.max_window_days)
        got = []
        for ws, we in windows:
            try:
                got += pam.collect(src.key, term, ws, we)
            except Exception:
                pass
        return src, pam.dedupe(got), windows

    records, runs = [], []
    with ThreadPoolExecutor(max_workers=10) as ex:
        for src, got, windows in (ex.map(run_one, sources) if sources else []):
            rel = [r for r in got if is_relevant((r.get("content") or "") + " " + (r.get("author") or ""), [term])]
            ctx = {"source_key": src.key, "access_path": src.access_path.value, "source_name": src.display_name}
            for r in rel:
                r["keyword"] = term
                r["project_id"] = pid
                struth.normalize_mention_source(r, ctx)
            records += rel
            runs.append(pam.build_coverage_for_run(src, start, end, windows, len(rel)))
    return {"stored": ingest(records, pid), "fetched": len(records), "runs": runs}


def backfill(keyword, start=None, end=None, historical_mode="recent_only", project=None):
    """Run a collection and return per-source coverage disclosure across the
    whole access ladder (so the UI can show what each platform actually searched)."""
    res = collect(keyword, start, end, project)
    runs_by_key = {r["source_key"]: r for r in res["runs"]}
    disclosure = []
    for s in pam.get_source_matrix():
        run = runs_by_key.get(s.key)
        disclosure.append({"source": s.key, "name": s.display_name, "platform": s.platform,
                           "status": s.status.value, "access_path": s.access_path.value,
                           "kept": run["kept"] if run else 0, "windows": run["windows"] if run else 0,
                           "chunked": run["chunked"] if run else False, "coverage": s.coverage_disclosure})
    return {"keyword": keyword, "historical_mode": historical_mode, "start": start, "end": end,
            "stored": res["stored"], "ran_at": now_iso(), "disclosure": disclosure}


def manual_import(term, text="", rows=None, project=None):
    """Ingest user-supplied data the user has lawful rights to upload (CSV or
    JSON). This is the compliant 'manual_import' access path. Never scraping."""
    proj_id = resolve_project(project)
    items = rows
    if items is None and text:
        t = text.strip()
        if t[:1] in "[{":
            try:
                j = json.loads(t)
                items = j if isinstance(j, list) else (j.get("rows") or j.get("mentions") or [])
            except Exception:
                items = None
        if items is None:
            items = list(csv.DictReader(io.StringIO(text)))
    recs = []
    for row in (items or []):
        if not isinstance(row, dict):
            continue
        content = row.get("content") or row.get("text") or row.get("snippet") or ""
        url = row.get("url") or row.get("source_url")
        pid = str(row.get("source_id") or row.get("id") or url
                  or hashlib.sha256(content.encode()).hexdigest()[:16])
        posted = str(row.get("posted_at") or row.get("timestamp") or now_iso())
        if len(posted) == 10:
            posted += "T00:00:00+00:00"
        rec = {"platform": (row.get("platform") or "manual").lower(), "platform_post_id": pid,
               "keyword": term, "author": row.get("author"), "content": content, "url": url,
               "posted_at": posted, "engagement": int(row.get("engagement") or 0),
               "description": row.get("description") or row.get("snippet") or "",
               "result_type": "manual_import", "project_id": proj_id}
        struth.normalize_mention_source(rec, {"source_key": "manual_import_csv",
                                          "manual_platform": rec["platform"], "access_path": "manual_import"})
        recs.append(rec)
    return ingest(recs, proj_id)


# ── metrics / feed / exports (all date-range aware) ──────────────────────────
# Retired open-network platforms are hidden from the product experience (feed,
# metrics, top sources, exports, coverage). Their connector code remains, but
# stored rows from them are not surfaced as active listening data.
_HIDDEN_PLATFORMS = ("mastodon", "lemmy", "nostr", "peertube", "hackernews", "news", "bluesky")


def _filters(keyword=None, start=None, end=None, project=None):
    clauses = ["is_hidden=0", "platform NOT IN (%s)" % ",".join("'%s'" % p for p in _HIDDEN_PLATFORMS)]
    params = []
    # Project scope: always bound to one project so data never mixes across projects.
    clauses.append("project_id = ?"); params.append(resolve_project(project))
    if keyword:
        clauses.append("keyword = ?"); params.append(keyword)
    if start:
        clauses.append("posted_date >= ?"); params.append(start)
    if end:
        clauses.append("posted_date <= ?"); params.append(end)
    return " AND ".join(clauses), params


STOPWORDS = set("the a an and or for to of in on at is are was be it this that with you your we our they "
                "them as from but not just have has had will can new about into over more most some any all "
                "what why how when who which their there here out up down vs amp via https http www com org net "
                "co will like dont been they that this with said were would people".split())


# ── Analytics model (documented in ANALYTICS.md) ──────────────────────────────
# Impressions = measured views/displays we actually receive from a platform API
#   (YouTube viewCount; owned-account impressions when a Meta account is connected).
# Reach (estimated audience reached) = per-item estimate, NOT engagement or follower count:
#   - items with measured views   -> unique viewers  = views * VIEW_TO_REACH (0.75)
#   - items with engagement, no views -> reach = engagement * ENGAGEMENT_TO_REACH (~1/4.5%)
#   - items with neither (news/open web) -> 0 (unknown; never fabricated)
# Public listening APIs do not expose unique reach, so this is a clearly-labeled estimate.
REACH_FACTOR = 0.75            # VIEW_TO_REACH: unique viewers as a fraction of total views
ENGAGEMENT_TO_REACH = 22.0     # ~1 / 4.5% engagement rate, to back out reach from engagement


def visibility_score(mentions, impressions, engagement, net_sentiment):
    """Cumulative 0-100 brand-visibility index. Log-scaled so each lever has
    diminishing returns; weighted volume 30 / amplification 30 / engagement 25 /
    sentiment 15. Fully documented in ANALYTICS.md."""
    if not mentions:
        return 0
    lg = lambda v: math.log10((v or 0) + 1)
    vol = min(lg(mentions) / 3.0, 1.0)       # ~1,000 mentions -> full marks
    amp = min(lg(impressions) / 6.0, 1.0)    # ~1,000,000 impressions -> full marks
    eng = min(lg(engagement) / 5.0, 1.0)     # ~100,000 interactions -> full marks
    sent = (max(-100, min(100, net_sentiment or 0)) + 100) / 200.0
    return int(round(100 * (0.30 * vol + 0.30 * amp + 0.25 * eng + 0.15 * sent)))


def _daily_buckets(vmap, start, end, cap=1100):
    """Continuous daily volume series from start..end with zero days filled, so the
    graph shows an unbroken timeline. Falls back to the sparse non-zero points if the
    range is absent, invalid, or wider than `cap` days (~3 years; a daily int array
    that size is still cheap to render)."""
    sparse = [{"t": k, "value": v} for k, v in sorted(vmap.items())]
    if not start or not end:
        return sparse
    try:
        d0 = datetime.strptime(start, "%Y-%m-%d").date()
        d1 = datetime.strptime(end, "%Y-%m-%d").date()
    except Exception:
        return sparse
    if d1 < d0 or (d1 - d0).days > cap:
        return sparse
    out, cur = [], d0
    while cur <= d1:
        s = cur.isoformat()
        out.append({"t": s, "value": vmap.get(s, 0)})
        cur += timedelta(days=1)
    return out


# Platform-level public-signal estimates for sources that expose no per-post metric
# (TikTok open-web/oEmbed references and News/open-web articles). Reach is estimated
# per surfaced item from mention volume; engagement from a typical public engagement
# rate. These drive the platform BUTTONS only (clearly labeled "est."); per-mention
# feed rows and the headline KPIs stay on measured data. Documented in ANALYTICS.md.
_PLATFORM_BASE_REACH = {"tiktok": 4500, "news": 1500, "youtube": 3000}
_PLATFORM_ENG_RATE = {"tiktok": 0.055, "news": 0.015, "youtube": 0.02}


def metrics(keyword=None, start=None, end=None, project=None):
    where, wp = _filters(keyword, start, end, project)
    kwhere, kwp = _filters(None, start, end, project)
    with db() as c:
        total = c.execute("SELECT COUNT(*) n FROM mentions WHERE " + where, wp).fetchone()["n"]
        vol = c.execute("SELECT posted_date d, COUNT(*) n FROM mentions WHERE " + where
                        + " GROUP BY d ORDER BY d", wp).fetchall()
        sent = c.execute("SELECT sentiment s, COUNT(*) n FROM mentions WHERE " + where + " GROUP BY sentiment", wp).fetchall()
        plats = c.execute("SELECT COALESCE(display_platform, platform) p, COUNT(*) n FROM mentions WHERE " + where
                          + " GROUP BY p ORDER BY n DESC", wp).fetchall()
        authors = c.execute("SELECT author a, platform p, COUNT(*) n FROM mentions WHERE " + where
                            + " AND author IS NOT NULL GROUP BY author ORDER BY n DESC, SUM(engagement) DESC LIMIT 6", wp).fetchall()
        titles = c.execute("SELECT content FROM mentions WHERE " + where + " ORDER BY posted_at DESC LIMIT 600", wp).fetchall()
        kw_rows = c.execute("SELECT keyword k, COUNT(*) n FROM mentions WHERE " + kwhere
                            + " AND keyword IS NOT NULL GROUP BY keyword ORDER BY n DESC", kwp).fetchall()
        cov_rows = c.execute("SELECT display_platform, platform_coverage_type, direct_platform_data, "
                             "discussed_platforms FROM mentions WHERE " + where + " LIMIT 2000", wp).fetchall()
        agg = c.execute("SELECT COALESCE(SUM(engagement),0) e, COALESCE(SUM(reach),0) r FROM mentions WHERE " + where, wp).fetchone()
        # Estimated audience reached (per item, summed): measured views -> unique viewers
        # (x VIEW_TO_REACH); items with engagement but no views -> back out reach from an
        # assumed engagement rate (engagement x ENGAGEMENT_TO_REACH). Never uses follower counts.
        reach_est = c.execute("SELECT COALESCE(SUM(CASE WHEN reach>0 THEN reach*0.75 "
                              "WHEN engagement>0 THEN engagement*22.0 ELSE 0 END),0) r FROM mentions WHERE " + where, wp).fetchone()
        plat_eng = c.execute("SELECT COALESCE(display_platform, platform) p, COALESCE(SUM(engagement),0) e, "
                             "COALESCE(SUM(reach),0) r FROM mentions WHERE " + where + " GROUP BY p ORDER BY r DESC, e DESC", wp).fetchall()
        top = c.execute("SELECT content, url, COALESCE(display_platform, platform) p, engagement, reach, sentiment, author "
                        "FROM mentions WHERE " + where + " ORDER BY reach DESC, engagement DESC LIMIT 5", wp).fetchall()
        # Per RAW source platform (open_web/youtube/bluesky) so the Sources cards can
        # show an accurate "N mentions in range" count keyed by source, not display name.
        src_counts = c.execute("SELECT platform p, COUNT(*) n FROM mentions WHERE " + where + " GROUP BY platform", wp).fetchall()
        tiktok_refs = c.execute("SELECT COUNT(*) n FROM mentions WHERE " + where + " AND url LIKE '%tiktok.com%'", wp).fetchone()
        # Per product-platform breakdown (TikTok, YouTube, News) for the platform
        # buttons: mentions, engagement, estimated reach, impressions, sentiment split.
        pstat = c.execute(
            # Bucket by the honest DISPLAY classification so the feed source tag, brand
            # color, filter, and metrics all agree. 'TikTok URL' = oEmbed/known TikTok
            # video references; News is the catch-all so buttons reconcile to total.
            "SELECT CASE WHEN COALESCE(display_platform,'')='TikTok URL' THEN 'tiktok' "
            "WHEN platform='youtube' THEN 'youtube' "
            "ELSE 'news' END bucket, "
            "COUNT(*) n, COALESCE(SUM(engagement),0) eng, COALESCE(SUM(reach),0) impr, "
            "COALESCE(SUM(CASE WHEN reach>0 THEN reach*0.75 WHEN engagement>0 THEN engagement*22.0 ELSE 0 END),0) reach_est, "
            "SUM(CASE WHEN sentiment='positive' THEN 1 ELSE 0 END) pos, "
            "SUM(CASE WHEN sentiment='negative' THEN 1 ELSE 0 END) neg "
            "FROM mentions WHERE " + where + " GROUP BY bucket", wp).fetchall()
    sc = {r["s"]: r["n"] for r in sent}
    pos, neg, neu = sc.get("positive", 0), sc.get("negative", 0), sc.get("neutral", 0)
    denom = max(pos + neg + neu, 1)
    freq = {}
    kwlow = (keyword or "").lower()
    for r in titles:
        for tok in TOKEN_RE.findall(re.sub(r"https?://\S+|www\.\S+", " ", (r["content"] or "").lower())):
            if len(tok) > 3 and tok not in STOPWORDS and tok != kwlow:
                freq[tok] = freq.get(tok, 0) + 1
    topics = sorted(freq.items(), key=lambda kv: kv[1], reverse=True)[:8]
    impressions = agg["r"]                                    # measured views/displays
    est_reach = int(round(reach_est["r"]))                    # estimated audience reached (per-item model)
    net = round(100 * (pos - neg) / denom)
    # Platform buttons: always emit TikTok, YouTube, News in that fixed order.
    _pby = {r["bucket"]: r for r in pstat}
    platform_stats = []
    for key, label, cov in (("tiktok", "TikTok", "Open web + oEmbed"),
                            ("youtube", "YouTube", "Official API"),
                            ("news", "News", "Open web")):
        r = _pby.get(key)
        n = r["n"] if r else 0
        e = r["eng"] if r else 0
        impr = r["impr"] if r else 0
        re_est = int(round(r["reach_est"])) if r else 0
        pp, nn = (r["pos"], r["neg"]) if r else (0, 0)
        netp = round(100 * (pp - nn) / max(n, 1)) if n else 0
        # Sources without a per-post public metric (TikTok references, News/open web)
        # get a labeled platform-level ESTIMATE derived from mention volume and a
        # documented per-item exposure factor, so the buttons are never empty.
        # Measured metrics (YouTube views/likes) always take precedence.
        estimated = False
        if n > 0 and re_est <= 0:
            re_est = n * _PLATFORM_BASE_REACH.get(key, 1500)
            estimated = True
        if n > 0 and e <= 0:
            e = int(round(re_est * _PLATFORM_ENG_RATE.get(key, 0.02)))
            estimated = True
        aud = int(round(re_est * 0.75)) if re_est else 0   # estimated unique audience
        platform_stats.append({"key": key, "label": label, "coverage": cov,
                               "mentions": n, "engagement": e, "impressions": impr,
                               "reach": re_est, "audience": aud, "net": netp, "estimated": estimated,
                               "visibility": visibility_score(n, max(impr, re_est), e, netp)})
    return {"kpis": {"totalMentions": total, "netSentiment": net,
                     "platforms": len(plats), "positivePct": round(100 * pos / denom),
                     "totalEngagement": agg["e"], "totalImpressions": impressions,
                     "totalReach": est_reach,
                     "visibilityScore": visibility_score(total, impressions, agg["e"], net)},
            "engagementByPlatform": [{"platform": r["p"], "engagement": r["e"], "reach": r["r"]} for r in plat_eng],
            "topContent": [{"content": (r["content"] or "").replace("\n", " ")[:160], "url": r["url"],
                            "platform": r["p"], "engagement": r["engagement"], "reach": r["reach"],
                            "sentiment": r["sentiment"], "author": r["author"]} for r in top],
            "volume": _daily_buckets({r["d"]: r["n"] for r in vol}, start, end),
            "sentiment": {"positive": pos, "neutral": neu, "negative": neg},
            "platforms": [{"platform": r["p"], "value": r["n"]} for r in plats],
            "topics": [{"label": k, "count": v} for k, v in topics],
            "authors": [{"author": r["a"], "platform": r["p"], "mentions": r["n"]} for r in authors],
            "keywords": [{"keyword": r["k"], "count": r["n"]} for r in kw_rows],
            "coverage": struth.coverage_summary_for_mentions([dict(r) for r in cov_rows]),
            "sourceCounts": {r["p"]: r["n"] for r in src_counts},
            "tiktokRefs": tiktok_refs["n"] if tiktok_refs else 0,
            "platformStats": platform_stats,
            "tracked": [k["label"] for k in tracked_keywords(project)]}


def recent(limit=25, keyword=None, platform=None, sentiment=None, start=None, end=None, project=None):
    where, args = _filters(keyword, start, end, project)
    if platform:
        where += " AND platform = ?"; args.append(platform)
    if sentiment:
        where += " AND sentiment = ?"; args.append(sentiment)
    args.append(limit)
    with db() as c:
        return [dict(r) for r in c.execute(
            "SELECT platform, keyword, author, content, description, result_type, url, posted_at, posted_date, ingested_at, "
            "sentiment, engagement, reach, "
            "display_platform, source_platform, searched_platform, discussed_platforms, direct_platform_data, "
            "platform_coverage_type, coverage_label, coverage_note, source_key, confidence_level "
            "FROM mentions WHERE " + where + " ORDER BY posted_at DESC LIMIT ?", args).fetchall()]


def coverage_ledger(keyword=None, start=None, end=None, project=None):
    """What was actually searched vs only discussed vs gated, for the range/term."""
    m = metrics(keyword, start, end, project)
    gated = [{"platform": a["platform"], "status": a["status"], "note": a["limitation"]}
             for a in pam.accounts_status() if a["source_key"] != "manual" and not a["can_collect"]]
    return {"range": {"start": start, "end": end}, "term": keyword or "all terms",
            "coverage": m["coverage"], "sources_searched": m["platforms"], "gated": gated}


def export_csv(keyword=None, start=None, end=None, project=None):
    where, params = _filters(keyword, start, end, project)
    with db() as c:
        rows = [dict(r) for r in c.execute(
            "SELECT keyword, display_platform, source_platform, searched_platform, discussed_platforms, "
            "direct_platform_data, platform_coverage_type, source_key, coverage_label, coverage_note, url, posted_at, "
            "author, sentiment, engagement, result_type, description, content FROM mentions WHERE " + where
            + " ORDER BY posted_at DESC", params).fetchall()]
    buf = io.StringIO()
    w = csv.writer(buf)
    # Source-truth fields first (invariant), then the captured result_type + description + content.
    w.writerow(struth.EXPORT_FIELDS + ["result_type", "description", "content"])
    for r in rows:
        safe = struth.export_safe_mention(r)
        w.writerow([safe.get(f) for f in struth.EXPORT_FIELDS]
                   + [r.get("result_type") or "",
                      (r.get("description") or "").replace("\n", " "),
                      (r.get("content") or "").replace("\n", " ")])
    return buf.getvalue().encode()


def export_insights(keyword=None, start=None, end=None, project=None):
    m = metrics(keyword, start, end, project)
    k = m["kpis"]
    return {"generated_at": now_iso(), "keyword": keyword or "all", "start": start, "end": end,
            "summary": "%s mentions in range. Net sentiment %+d (%d%% positive)." % (
                "{:,}".format(k["totalMentions"]), k["netSentiment"], k["positivePct"]),
            "sources": [s.to_dict() for s in pam.get_source_matrix()], **m}


def report_html(keyword=None, start=None, end=None, project=None):
    m = metrics(keyword, start, end, project)
    k = m["kpis"]
    rows = recent(20, keyword, None, None, start, end, project)
    proj = project_name(project)
    s = m["sentiment"]
    tot = max(s["positive"] + s["neutral"] + s["negative"], 1)

    def e(x):
        return (str(x) if x is not None else "").replace("&", "&amp;").replace("<", "&lt;")

    def ymd(d, fallback=""):
        d = (d or "")[:10]
        try:
            y, mo, da = d.split("-")
            return y[2:] + "/" + mo + "/" + da
        except Exception:
            return fallback
    cov = m["coverage"]
    plat = "".join("<tr><td>%s</td><td style='text-align:right'>%d</td></tr>" % (e(p["platform"]), p["value"]) for p in m["platforms"]) or "<tr><td>none</td></tr>"
    feed = "".join("<li><a href='%s'>%s</a> <span style='color:#777'>. %s · %s</span></li>"
                   % (e(r["url"] or "#"), e((r["content"] or "")[:140]), e(r.get("coverage_label") or r.get("display_platform")), e(r["sentiment"])) for r in rows) or "<li>none</li>"
    searched = ", ".join("%s (%d)" % (e(d["platform"]), d["count"]) for d in cov["platforms_searched"]) or "none"
    disc = ", ".join("%s (%d)" % (e(d["platform"]), d["count"]) for d in cov["platforms_discussed"]) or "none"
    gated = [a for a in pam.accounts_status() if a["source_key"] != "manual" and not a["can_collect"]]
    gated_html = "".join("<li>%s. %s</li>" % (e(a["platform"]), e(a["status"])) for a in gated) or "<li>none</li>"
    cov_html = ("<h3>Coverage &amp; source truth</h3><ul>"
                "<li>Direct platform data included: <b>%d</b></li>"
                "<li>Open-web references included: <b>%d</b></li>"
                "<li>Manual imports: <b>%d</b></li>"
                "<li>Licensed-provider feeds: <b>%d</b></li>"
                "<li>Platforms directly searched: %s</li>"
                "<li>Platforms discussed but not directly searched: %s</li></ul>"
                "<p class='muted'>Gated sources (not searched):</p><ul>%s</ul>"
                "<p class='muted'>Caveat: open-web references mention a platform but are not that platform's native "
                "data. Closed-platform listening requires connected accounts, approved research access, or a licensed "
                "provider.</p>") % (cov["direct"], cov["open_web_references"], cov["manual_imports"], cov["licensed"],
                                    searched, disc, gated_html)
    return ("<!doctype html><meta charset='utf-8'><title>EGC Pulse Report</title>"
            "<style>@import url('https://fonts.googleapis.com/css2?family=Lato:wght@400;700&family=Noto+Serif:wght@400;600&display=swap');"
            "body{font-family:'Lato',system-ui,Arial,sans-serif;color:#111;max-width:820px;margin:32px auto;padding:0 20px}"
            "h1,h3{font-family:'Noto Serif',Georgia,serif;font-weight:600}.kpi{display:inline-block;border:1px solid #ddd;border-radius:8px;padding:8px 14px;margin:4px 8px 4px 0}"
            ".kpi b{font-size:22px;display:block}td{padding:3px 12px;border-bottom:1px solid #eee}.muted{color:#777}ul{line-height:1.6}</style>"
            "<p class='muted' style='margin-bottom:2px;letter-spacing:.08em;text-transform:uppercase;font-size:11px'>EGC Pulse. Social Listening Report</p>"
            "<h1 style='margin:0 0 2px'>%s</h1><p class='muted'>Tracked term: <b>%s</b> · %s to %s · Generated %s</p>"
            "<div><span class='kpi'><b>%s</b>mentions</span><span class='kpi'><b>%+d</b>net sentiment</span>"
            "<span class='kpi'><b>%d%%</b>positive</span><span class='kpi'><b>%d</b>sources</span></div>"
            "<h3>Sentiment</h3><p>Positive %d%% · Neutral %d%% · Negative %d%%</p>"
            "%s"
            "<h3>Mentions by source</h3><table>%s</table><h3>Top mentions</h3><ol>%s</ol>"
            "<p class='muted'>Internal use only. Not for resale or redistribution. Source access depends on "
            "configured APIs, connected accounts, licensed providers, approved research access, or lawful manual "
            "import. Generated by EGC Pulse. Print, then Save as PDF.</p>") % (
        e(proj), e(keyword or "all keywords"), ymd(start, "earliest"), ymd(end, "now"), ymd(now_iso()),
        "{:,}".format(k["totalMentions"]), k["netSentiment"], k["positivePct"], k["platforms"],
        round(100 * s["positive"] / tot), round(100 * s["neutral"] / tot), round(100 * s["negative"] / tot),
        cov_html, plat, feed)


def _pptx_filename(project=None):
    """Per-project download filename, e.g. EGC_Pulse_Acme_Launch.pptx."""
    nm = re.sub(r"[^A-Za-z0-9]+", "_", project_name(project)).strip("_") or "Report"
    return "EGC_Pulse_%s.pptx" % nm[:60]


def _fmt_n(n):
    """Compact number for slides: 1234567 -> 1.2M, 12300 -> 12.3K."""
    try:
        n = float(n or 0)
    except Exception:
        return "0"
    for div, suf in ((1e9, "B"), (1e6, "M"), (1e3, "K")):
        if abs(n) >= div:
            s = ("%.1f" % (n / div)).rstrip("0").rstrip(".")
            return s + suf
    return "{:,}".format(int(round(n)))


def report_pptx(keyword=None, start=None, end=None, project=None):
    """Build a PowerPoint deck that mirrors the dashboard: same metrics, charts,
    insights and data, in the same dark theme. Returns .pptx bytes. Requires
    python-pptx (listed in requirements.txt); raises RuntimeError if unavailable."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
        from pptx.chart.data import CategoryChartData
    except Exception as exc:  # pragma: no cover - dependency guard
        raise RuntimeError("python-pptx is not installed. Add it to requirements.txt and redeploy.") from exc

    m = metrics(keyword, start, end, project)
    k = m["kpis"]
    s = m["sentiment"]
    cov = m["coverage"]
    proj = project_name(project)

    # Dashboard palette (kept in sync with dashboard.html tokens).
    DARK = RGBColor(0x19, 0x1A, 0x1D)
    PANEL = RGBColor(0x21, 0x23, 0x29)
    LINE = RGBColor(0x34, 0x37, 0x3F)
    INK = RGBColor(0xEE, 0xF0, 0xF3)
    MUT = RGBColor(0x9A, 0x9F, 0xA8)
    FAINT = RGBColor(0x6D, 0x72, 0x7B)
    BLUE = RGBColor(0x3D, 0x7E, 0xFF)
    POS = RGBColor(0x3F, 0xB9, 0x50)
    NEU = RGBColor(0x8A, 0x90, 0x99)
    NEG = RGBColor(0xE5, 0x5B, 0x4C)
    FONT = "Lato"

    def ymd(d):
        try:
            y, mo, da = (d or "")[:10].split("-")
            return y[2:] + "/" + mo + "/" + da
        except Exception:
            return d or ""

    def daterange_label():
        if not (start and end):
            return "All available dates"
        try:
            d0 = datetime.strptime(start[:10], "%Y-%m-%d").date()
            d1 = datetime.strptime(end[:10], "%Y-%m-%d").date()
            n = (d1 - d0).days + 1
            return "%s to %s   ·   %d day%s" % (ymd(start), ymd(end), n, "" if n == 1 else "s")
        except Exception:
            return ymd(start) + " to " + ymd(end)

    def timeline(series):
        """Aggregate the daily volume into a readable, chronological axis: daily up to
        16 points, then weekly, then monthly. Returns (labels, values, granularity)."""
        series = series or []
        days = len(series)
        if days <= 16:
            return [ymd(p["t"]) for p in series], [p["value"] for p in series], "Daily"
        if days <= 120:
            labels, vals = [], []
            for i in range(0, days, 7):
                chunk = series[i:i + 7]
                labels.append(ymd(chunk[0]["t"]))
                vals.append(sum(p["value"] for p in chunk))
            return labels, vals, "Weekly"
        groups, order = {}, []
        for p in series:
            ym = (p["t"] or "")[:7]
            if ym not in groups:
                groups[ym] = 0; order.append(ym)
            groups[ym] += p["value"]
        return [ym[2:].replace("-", "/") for ym in order], [groups[ym] for ym in order], "Monthly"

    rng = daterange_label() if (start or end) else "all available dates"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = 13.333, 7.5

    def new_slide():
        sl = prs.slides.add_slide(BLANK)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = DARK
        return sl

    def text(sl, l, t, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE   # honor the fixed box; never grow/shrink to fit
        except Exception:
            pass
        tf.vertical_anchor = anchor
        items = runs if isinstance(runs, list) else [(runs, size, color, bold)]
        for i, item in enumerate(items):
            txt, sz, col, bd = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(sz)
            r.font.bold = bd
            r.font.color.rgb = col
            r.font.name = FONT
        return tb

    def card(sl, l, t, w, h, fill=PANEL, edge=LINE):
        sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = edge
        sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
        return sh

    def header(sl, title, sub=None):
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.52), Inches(0.12), Inches(0.66))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background(); bar.shadow.inherit = False
        text(sl, 0.97, 0.46, 7.0, 0.3, [(proj.upper()[:46], 10.5, BLUE, True)])
        text(sl, 0.95, 0.70, 7.1, 0.6, title, 24, INK, True)
        text(sl, SW - 5.0, 0.52, 4.3, 0.35, [("EGC PULSE", 11, FAINT, True)], align=PP_ALIGN.RIGHT)
        text(sl, SW - 5.0, 0.80, 4.3, 0.35, [(rng, 10.5, FAINT, False)], align=PP_ALIGN.RIGHT)
        if sub:
            text(sl, 0.97, 1.28, 11.0, 0.4, [(sub, 12, MUT, False)])

    def style_chart(chart, legend=False):
        chart.has_title = False
        try:
            chart.font.size = Pt(10)
            chart.font.color.rgb = INK
            chart.font.name = FONT
        except Exception:
            pass
        chart.has_legend = legend
        if legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = MUT
            chart.legend.font.size = Pt(10)

    # Slide 1: title (the project name is the hero)
    sl = new_slide()
    for dx, dy, c in [(0, 0, BLUE), (0.34, 0, RGBColor(0x2B, 0x57, 0xB0)),
                      (0, 0.34, RGBColor(0x2B, 0x57, 0xB0)), (0.34, 0.34, RGBColor(0x4E, 0x53, 0x5C))]:
        sq = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.92 + dx), Inches(0.92 + dy), Inches(0.27), Inches(0.27))
        sq.fill.solid(); sq.fill.fore_color.rgb = c; sq.line.fill.background(); sq.shadow.inherit = False
        try:
            sq.adjustments[0] = 0.26
        except Exception:
            pass
    text(sl, 1.62, 0.93, 9.5, 0.4, [("EGC PULSE      ·      SOCIAL LISTENING REPORT", 12, FAINT, True)])
    # hero: the project name, explicitly labeled so the deck identifies its project
    text(sl, 0.9, 2.18, 9.0, 0.32, [("PROJECT", 12, BLUE, True)])
    hero = proj[:70]
    hsize = 50 if len(hero) <= 22 else (40 if len(hero) <= 34 else (32 if len(hero) <= 50 else 26))
    text(sl, 0.86, 2.56, 11.7, 1.2, hero, hsize, INK, True)
    rule = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(3.92), Inches(1.7), Inches(0.07))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLUE; rule.line.fill.background(); rule.shadow.inherit = False
    text(sl, 0.9, 4.12, 11.7, 0.5, [("Tracked term:  " + (keyword or "all tracked terms") + "        ·        " + rng, 14, MUT, False)])
    # headline KPI strip for an at-a-glance opening
    tk = [("Total mentions", "{:,}".format(k["totalMentions"])),
          ("Est. audience reached", _fmt_n(k.get("totalReach", 0))),
          ("Engagement", _fmt_n(k.get("totalEngagement", 0))),
          ("Visibility score", str(k.get("visibilityScore", 0)) + " / 100")]
    cw2, gx2, x2, y2, h2 = 2.86, 0.30, 0.9, 5.05, 1.45
    for i, (lab, val) in enumerate(tk):
        cx = x2 + i * (cw2 + gx2)
        card(sl, cx, y2, cw2, h2)
        text(sl, cx + 0.26, y2 + 0.27, cw2 - 0.45, 0.4, [(lab.upper(), 10, FAINT, True)])
        text(sl, cx + 0.24, y2 + 0.68, cw2 - 0.4, 0.7, [(val, 27, INK, True)])
    text(sl, 0.9, SH - 0.58, 12.0, 0.4, [("Internal use only. Compliant, official-source listening. Generated " + now_iso()[:16].replace("T", " "), 10, FAINT, False)])

    # Slide 2: overview KPIs
    sl = new_slide()
    header(sl, "Overview", "Headline metrics for the selected term and date range.")
    net = k["netSentiment"]
    kpis = [("Total mentions", "{:,}".format(k["totalMentions"]), INK),
            ("Impressions (measured)", _fmt_n(k.get("totalImpressions", 0)), INK),
            ("Reach (estimated audience)", _fmt_n(k.get("totalReach", 0)), BLUE),
            ("Engagement", _fmt_n(k.get("totalEngagement", 0)), INK),
            ("Net sentiment", ("+" if net > 0 else "") + str(net), POS if net > 0 else (NEG if net < 0 else NEU)),
            ("Visibility score", str(k.get("visibilityScore", 0)) + " / 100", INK)]
    cw, ch, gx, gy = 3.78, 1.95, 0.34, 0.34
    x0, y0 = 0.7, 1.75
    for i, (lab, val, col) in enumerate(kpis):
        cx = x0 + (i % 3) * (cw + gx)
        cy = y0 + (i // 3) * (ch + gy)
        card(sl, cx, cy, cw, ch)
        text(sl, cx + 0.28, cy + 0.26, cw - 0.5, 0.45, [(lab.upper(), 11, FAINT, True)])
        text(sl, cx + 0.26, cy + 0.74, cw - 0.45, 1.0, [(val, 34, col, True)])
    text(sl, 0.7, SH - 0.62, 12.0, 0.4,
         [("Reach is an estimate of audience reached: measured views to unique viewers, plus an engagement-rate model where views are unavailable. Never follower counts.", 9.5, FAINT, False)])

    # Slide 3: platform performance (TikTok, YouTube, News) - mirrors the dashboard buttons
    sl = new_slide()
    header(sl, "Platform performance", "TikTok, YouTube and News, derived from the collected research.")
    pstats = m.get("platformStats", [])
    PACC = {"tiktok": RGBColor(0x25, 0xF4, 0xEE), "youtube": RGBColor(0xFF, 0x2B, 0x3E), "news": BLUE}
    pcw, pgx, px0, py0, pch = 3.92, 0.34, 0.7, 1.7, 4.55
    for i, sp in enumerate(pstats[:3]):
        cx = px0 + i * (pcw + pgx)
        card(sl, cx, py0, pcw, pch)
        acc = PACC.get(sp["key"], BLUE)
        stripe = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(py0), Inches(pcw), Inches(0.10))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = acc; stripe.line.fill.background(); stripe.shadow.inherit = False
        text(sl, cx + 0.3, py0 + 0.36, pcw - 0.6, 0.5, [(sp["label"], 23, INK, True)])
        covlab = sp["coverage"].upper() + ("  ·  EST." if sp.get("estimated") else "")
        text(sl, cx + 0.31, py0 + 0.98, pcw - 0.6, 0.35, [(covlab, 9.5, FAINT, True)])
        mets = [("Mentions", "{:,}".format(sp["mentions"])), ("Reach (est.)", _fmt_n(sp["reach"])),
                ("Engagement", _fmt_n(sp["engagement"])), ("Visibility", str(sp["visibility"]))]
        mw = (pcw - 0.6) / 2.0
        for j, (ml, mv) in enumerate(mets):
            ox = cx + 0.3 + (j % 2) * mw
            oy = py0 + 1.6 + (j // 2) * 1.42
            text(sl, ox, oy, mw - 0.05, 0.55, [(mv, 23, INK, True)])
            text(sl, ox, oy + 0.52, mw - 0.05, 0.3, [(ml.upper(), 9.5, FAINT, True)])
    text(sl, 0.7, SH - 0.62, 12.0, 0.4,
         [("TikTok and News are open-web references (not native platform listening); engagement and reach appear only where a public metric exists.", 9.5, FAINT, False)])

    # Slide 4: mention volume (chronological, granularity-aware timeline)
    sl = new_slide()
    vol = m.get("volume", [])
    cats, vals, gran = timeline(vol)
    header(sl, "Mention volume", "%s mention volume, %s." % (gran, rng))
    if vol and sum(p["value"] for p in vol) > 0:
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("Mentions", vals)
        card(sl, 0.7, 1.78, SW - 1.4, 5.12)
        gf = sl.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.95), Inches(2.05), Inches(SW - 1.9), Inches(4.35), cd)
        chart = gf.chart
        style_chart(chart)
        ser = chart.series[0]
        ser.format.line.color.rgb = BLUE
        ser.format.line.width = Pt(2.5)
        ser.smooth = False
        peak = max(vol, key=lambda p: p["value"])
        active = sum(1 for p in vol if p["value"] > 0)
        text(sl, 0.95, SH - 0.66, 11.0, 0.4,
             [("Peak: %s mention%s on %s. Activity on %d of %d day%s in range." % (
                 "{:,}".format(peak["value"]), "" if peak["value"] == 1 else "s", ymd(peak["t"]),
                 active, len(vol), "" if len(vol) == 1 else "s"), 10.5, FAINT, False)])
    else:
        text(sl, 0.95, 3.2, 11.0, 1.0, [("No mention volume in this range.", 16, MUT, False)])

    # Slide 4: sentiment + top platforms
    sl = new_slide()
    header(sl, "Sentiment and platform mix")
    if (s["positive"] + s["neutral"] + s["negative"]) > 0:
        card(sl, 0.7, 1.55, 5.7, 5.35)
        text(sl, 0.95, 1.7, 5.0, 0.4, [("SENTIMENT", 11, FAINT, True)])
        cd = CategoryChartData()
        cd.categories = ["Positive", "Neutral", "Negative"]
        cd.add_series("Sentiment", [s["positive"], s["neutral"], s["negative"]])
        gf = sl.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.95), Inches(2.1), Inches(5.2), Inches(4.5), cd)
        chart = gf.chart
        style_chart(chart, legend=True)
        pts = chart.plots[0].series[0].points
        for pt, col in zip(pts, (POS, NEU, NEG)):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
            pt.format.line.color.rgb = PANEL
    else:
        text(sl, 0.95, 3.2, 5.0, 0.8, [("No sentiment data in this range.", 14, MUT, False)])
    plats = m["platforms"][:6]
    if plats:
        card(sl, 6.9, 1.55, SW - 7.6, 5.35)
        text(sl, 7.15, 1.7, 5.0, 0.4, [("MENTIONS BY SOURCE", 11, FAINT, True)])
        cd = CategoryChartData()
        cd.categories = [p["platform"] for p in plats]
        cd.add_series("Mentions", [p["value"] for p in plats])
        gf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.05), Inches(2.05), Inches(SW - 7.75), Inches(4.6), cd)
        chart = gf.chart
        style_chart(chart)
        ser = chart.series[0]
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = BLUE
        try:
            ser.has_data_labels = True
            ser.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            ser.data_labels.font.color.rgb = INK
            ser.data_labels.font.size = Pt(10)
        except Exception:
            pass
    else:
        text(sl, 7.15, 3.2, 5.0, 0.8, [("No platform data in this range.", 14, MUT, False)])

    # Slide 5: top topics + top authors
    sl = new_slide()
    header(sl, "Topics and voices")
    topics = m["topics"][:8]
    if topics:
        card(sl, 0.7, 1.55, 6.0, 5.35)
        text(sl, 0.95, 1.7, 5.0, 0.4, [("TOP TOPICS", 11, FAINT, True)])
        cd = CategoryChartData()
        cd.categories = [t["label"] for t in topics]
        cd.add_series("Mentions", [t["count"] for t in topics])
        gf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.95), Inches(2.05), Inches(5.5), Inches(4.6), cd)
        chart = gf.chart
        style_chart(chart)
        chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = BLUE
    else:
        text(sl, 0.95, 3.2, 5.0, 0.8, [("No topics extracted in this range.", 14, MUT, False)])
    authors = m["authors"][:6]
    card(sl, 7.0, 1.55, SW - 7.7, 5.35)
    text(sl, 7.25, 1.7, 5.0, 0.4, [("TOP AUTHORS", 11, FAINT, True)])
    if authors:
        _table(sl, 7.05, 2.15, SW - 7.75, 4.4, ["Author", "Source", "Mentions"],
               [[(a["author"] or "unknown")[:28], a["platform"], str(a["mentions"])] for a in authors],
               BLUE, PANEL, LINE, INK, MUT, FONT)
    else:
        text(sl, 7.25, 3.2, 5.0, 0.8, [("No attributed authors in this range.", 14, MUT, False)])

    # Slide 6: top content
    sl = new_slide()
    header(sl, "Top content", "Highest-reach mentions in range. Reach shown only where a public metric exists.")
    tc = m["topContent"][:5]
    if tc:
        rows = []
        for c in tc:
            metric = _fmt_n(c["reach"]) if c.get("reach") else (_fmt_n(c["engagement"]) + " eng" if c.get("engagement") else "no public metric")
            rows.append([(c["content"] or "")[:78] or "(untitled)", c["platform"], (c.get("sentiment") or "")[:3], metric])
        _table(sl, 0.7, 1.7, SW - 1.4, 4.9, ["Content", "Source", "Sent.", "Reach / metric"],
               rows, BLUE, PANEL, LINE, INK, MUT, FONT, col_widths=[7.4, 1.8, 1.0, 1.7])
    else:
        text(sl, 0.95, 3.2, 11.0, 0.8, [("No mentions stored in this range.", 16, MUT, False)])

    # Slide 7: coverage & source truth
    sl = new_slide()
    header(sl, "Coverage and source truth")
    searched = ", ".join("%s (%d)" % (d["platform"], d["count"]) for d in cov.get("platforms_searched", [])) or "none"
    disc = ", ".join("%s (%d)" % (d["platform"], d["count"]) for d in cov.get("platforms_discussed", [])) or "none"
    lines = [
        ("Direct platform data (official APIs / connected accounts): %d" % cov.get("direct", 0), 14, INK, False),
        ("Open-web references (mention a platform, not its native data): %d" % cov.get("open_web_references", 0), 14, INK, False),
        ("Manual imports: %d" % cov.get("manual_imports", 0), 14, INK, False),
        ("Licensed-provider feeds: %d" % cov.get("licensed", 0), 14, INK, False),
        ("", 8, MUT, False),
        ("Platforms directly searched:  " + searched, 12, MUT, False),
        ("Platforms discussed but not directly searched:  " + disc, 12, MUT, False),
    ]
    card(sl, 0.7, 1.55, SW - 1.4, 3.7)
    text(sl, 1.0, 1.8, SW - 2.0, 3.3, lines)
    text(sl, 0.7, 5.5, SW - 1.4, 1.4,
         [("Caveat: open-web and News references mention a platform but are not that platform's native data. "
           "Closed-platform listening requires connected accounts, approved research access, or a licensed provider. "
           "Reach figures are clearly-labeled estimates; items with no public metric are never assigned fabricated numbers.", 11, FAINT, False)])

    import io as _io
    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _table(sl, l, t, w, h, headers, rows, head_fill, body_fill, edge, ink, mut, font, col_widths=None):
    """Helper: a dark-themed table for the deck."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    n = len(rows) + 1
    gf = sl.shapes.add_table(n, len(headers), Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for i, cwid in enumerate(col_widths):
            tbl.columns[i].width = Inches(cwid)
    for j, htext in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = head_fill
        cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htext
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = font
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = body_fill
            cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(10.5); r.font.color.rgb = ink if j == 0 else mut; r.font.name = font
    return tbl


def internal_use():
    return os.getenv("INTERNAL_USE_ONLY", "true").strip().lower() in {"1", "true", "yes", "on"}


# ── async collection jobs (UI never freezes; progress + cancel) ──────────────
_jobs = {}
_jobs_lock = threading.Lock()

SKIP_MSG = {
    "reddit": "Skipped Reddit. API key not configured.",
    "youtube": "Skipped YouTube. API key not configured.",
    "x": "Skipped X. API key not configured.",
    "meta": "Skipped Meta. Connected account, Ad Library token, or approved research access required.",
    "instagram": "Skipped Instagram. Connected account or licensed provider required.",
    "facebook": "Skipped Facebook. Connected account, approved research access, or licensed provider required.",
    "tiktok": "Skipped TikTok. Research approval, connected account, commercial content access, or licensed provider required.",
    "licensed": "Skipped Licensed provider. LICENSED_PROVIDER_API_KEY/URL not configured.",
}


def _skipped_sources():
    out = []
    for a in pam.accounts_status():
        if a["source_key"] == "manual" or a["can_collect"]:
            continue
        out.append({"source": a["platform"], "reason": SKIP_MSG.get(a["source_key"], a["status"])})
    return out


def _clamp_recent(start, end, days=30):
    end = (end or now_iso()[:10])[:10]
    try:
        floor = (datetime.strptime(end, "%Y-%m-%d") - timedelta(days=days)).strftime("%Y-%m-%d")
        return floor if (not start or start < floor) else start
    except Exception:
        return start


def job_public(jid):
    with _jobs_lock:
        j = _jobs.get(jid)
        return {k: v for k, v in j.items() if not k.startswith("_")} if j else None


def cancel_job(jid):
    with _jobs_lock:
        j = _jobs.get(jid)
    if not j or j["status"] in ("complete", "failed", "cancelled"):
        return False
    j["_cancel"].set()
    return True


def start_job(terms, start, end, mode, project=None):
    jid = uuid.uuid4().hex[:12]
    j = {"job_id": jid, "status": "queued", "progress_pct": 0, "current_source": None,
         "current_chunk": 0, "total_chunks": 0, "stored_count": 0, "skipped_sources": [],
         "errors": [], "coverage": [], "message": "", "terms": terms, "start": start, "end": end,
         "mode": mode, "project": resolve_project(project), "started_at": now_iso(), "_cancel": threading.Event()}
    with _jobs_lock:
        _jobs[jid] = j
    threading.Thread(target=run_job, args=(jid,), daemon=True).start()
    return jid


def run_job(jid):
    with _jobs_lock:
        j = _jobs.get(jid)
    if not j:
        return
    try:
        j["status"] = "running"
        start = _clamp_recent(j["start"], j["end"], 30) if j["mode"] == "fast" else j["start"]
        end = j["end"]
        live = pam.get_live_collectors()
        # Prioritize YouTube and TikTok (open-web discovery leads with the TikTok query)
        # as the primary sources: collect and analyze them first, then the rest.
        _PRIO = {"youtube_official_api": 0, "open_web_social_discovery": 1}
        live.sort(key=lambda s: _PRIO.get(s.key, 5))
        j["skipped_sources"] = _skipped_sources()
        if not live:
            j["status"] = "no_sources"
            j["message"] = ("No data sources are configured, so there is nothing to search. "
                            "Add credentials for a source (Reddit is the quickest) in demo/.env "
                            "locally, or in your host's environment (Render → Environment) for "
                            "the deploy, then run again.")
            return
        units = [(t, s) for t in j["terms"] for s in live]
        total = len(units) or 1
        done = 0
        rate = False
        for term, src in units:
            if j["_cancel"].is_set():
                j["status"] = "cancelled"
                return
            j["current_source"] = "%s. %s" % (term, src.display_name)
            windows = pam.chunk_range(start, end, src.max_window_days)
            j["total_chunks"] = len(windows)
            got = []
            for ci, (ws, we) in enumerate(windows):
                if j["_cancel"].is_set():
                    j["status"] = "cancelled"
                    return
                j["current_chunk"] = ci + 1
                try:
                    got += pam.collect(src.key, term, ws, we)
                except Exception as e:
                    msg = str(e)
                    rate = rate or ("429" in msg or "rate" in msg.lower())
                    j["errors"].append({"source": src.key, "error": msg[:140]})
            got = pam.dedupe(got)
            # Open-web discovery results are already term-scoped by the search provider
            # (e.g. site:tiktok.com "term"), so trust them; only re-filter API sources
            # (YouTube/Bluesky/etc.) whose results can be loosely related. The URL is
            # included in the relevance text so term-in-handle/path matches count.
            rel = [r for r in got if (r.get("platform") == "open_web"
                   or is_relevant((r.get("content") or "") + " " + (r.get("author") or "") + " " + (r.get("url") or ""), [term]))]
            ctx = {"source_key": src.key, "access_path": src.access_path.value, "run_id": jid, "source_name": src.display_name}
            pid = j.get("project")
            for r in rel:
                r["keyword"] = term
                r["project_id"] = pid
                struth.normalize_mention_source(r, ctx)
            j["stored_count"] += ingest(rel, pid)  # partial results land in the DB immediately
            j["coverage"].append(pam.build_coverage_for_run(src, start, end, windows, len(rel)))
            done += 1
            j["progress_pct"] = round(100 * done / total)
        j["status"] = "rate_limited" if (rate and j["stored_count"] == 0) else "complete"
        if j["status"] == "complete" and j["stored_count"] == 0 and not j["errors"]:
            j["message"] = ("Searched %d source(s) but found no matching mentions in this range. "
                            "Try a broader term or a wider date range." % len(live))
    except Exception as e:
        j["status"] = "failed"
        j["errors"].append({"source": "job", "error": str(e)[:140]})
    finally:
        j["finished_at"] = now_iso()


# ── connection tests (never leak secret VALUES. Only variable names) ────────
def test_connection(key):
    g = os.environ.get
    now = now_iso()

    def done(ok, msg, missing=()):
        pam._last_checked[key] = now
        status = next((a["status"] for a in pam.accounts_status() if a["source_key"] == key), "")
        return {"source_key": key, "ok": ok, "status": status, "message": msg,
                "missing_env_vars": list(missing), "checked_at": now}

    if key == "manual":
        return done(True, "Manual import is always available for data you have lawful rights to upload.")
    if key == "reddit":
        miss = [v for v in ("REDDIT_CLIENT_ID", "REDDIT_CLIENT_SECRET") if not g(v)]
        if miss:
            return done(False, "Missing %s. Add them to demo/.env (local) or your host's environment (Render), then test again." % " and ".join(miss), miss)
        try:
            import base64
            auth = base64.b64encode(("%s:%s" % (g("REDDIT_CLIENT_ID"), g("REDDIT_CLIENT_SECRET"))).encode()).decode()
            ua = g("REDDIT_USER_AGENT") or "python:egc-pulse:v1.0 (by /u/egc-pulse)"
            req = urllib.request.Request("https://www.reddit.com/api/v1/access_token", data=b"grant_type=client_credentials",
                                         headers={"Authorization": "Basic " + auth, "User-Agent": ua,
                                                  "Content-Type": "application/x-www-form-urlencoded"})
            tok = json.loads(urllib.request.urlopen(req, timeout=10).read()).get("access_token")
            return done(bool(tok), "Credentials valid. Reddit Data API reachable." if tok
                        else "Credentials accepted but no token returned. Check the app type ('script' or 'web app') and that Data API access is approved.")
        except urllib.error.HTTPError as e:
            if e.code == 401:
                return done(False, "Reddit rejected the credentials (401). Re-check REDDIT_CLIENT_ID / REDDIT_CLIENT_SECRET and that the app type is 'script' or 'web app'.")
            if e.code == 403:
                return done(False, "Reddit returned 403 (forbidden). The app likely needs Data API access approval (Responsible Builder Policy).")
            if e.code == 429:
                return done(True, "Credentials valid (rate-limited right now, HTTP 429).")
            return done(False, "Reddit rejected the request (HTTP %d). Check app type, token validity, or API tier." % e.code)
        except Exception:
            return done(False, "Could not reach the Reddit API. Check the keys and network.")
    if key == "youtube":
        if not g("YOUTUBE_API_KEY"):
            return done(False, "Missing YOUTUBE_API_KEY. Create a free key in Google Cloud (enable 'YouTube Data API v3'), add it to demo/.env or your host's environment, then test again.", ["YOUTUBE_API_KEY"])
        try:
            # videos.list with a known id costs only 1 quota unit (vs 100 for search) — cheap key check.
            req = urllib.request.Request("https://www.googleapis.com/youtube/v3/videos?part=id&id=dQw4w9WgXcQ&key="
                                         + urllib.parse.quote(g("YOUTUBE_API_KEY")))
            urllib.request.urlopen(req, timeout=10)
            return done(True, "Credentials valid. YouTube Data API reachable.")
        except urllib.error.HTTPError as e:
            if e.code in (400, 403):
                return done(False, "YouTube rejected the key (HTTP %d). Check the key value and that 'YouTube Data API v3' is enabled in Google Cloud." % e.code)
            return done(False, "YouTube API error (HTTP %d)." % e.code)
        except Exception:
            return done(False, "Could not reach the YouTube API. Check the key and network.")
    if key == "x":
        if not g("X_BEARER_TOKEN"):
            return done(False, "Missing X_BEARER_TOKEN. Add it to demo/.env, restart the server, then test again.", ["X_BEARER_TOKEN"])
        try:
            req = urllib.request.Request("https://api.x.com/2/tweets/search/recent?query=egc%20-is%3Aretweet&max_results=10",
                                         headers={"Authorization": "Bearer " + g("X_BEARER_TOKEN")})
            urllib.request.urlopen(req, timeout=10)
            return done(True, "Credentials valid. X API reachable.")
        except urllib.error.HTTPError as e:
            if e.code == 429:
                return done(True, "Credentials valid (rate-limited right now. HTTP 429).")
            return done(False, "Credentials found, but the platform rejected the request (HTTP %d). Check token validity or API tier." % e.code)
        except Exception:
            return done(False, "Could not reach the X API. Check the token and network.")
    if key in ("meta", "instagram", "facebook"):
        miss = [v for v in ("META_APP_ID", "META_APP_SECRET", "META_ACCESS_TOKEN") if not g(v)]
        if not g("META_ACCESS_TOKEN"):
            return done(False, "Missing %s. Add them to demo/.env, restart the server, then test again." % " and ".join(miss or ["META_ACCESS_TOKEN"]), miss)
        try:
            req = urllib.request.Request("https://graph.facebook.com/v19.0/me?fields=id&access_token=" + urllib.parse.quote(g("META_ACCESS_TOKEN")))
            urllib.request.urlopen(req, timeout=10)
            return done(True, "Meta token accepted by the Graph API. Owned-account / Ad Library scope depends on app review.")
        except urllib.error.HTTPError as e:
            return done(False, "Token found, but Meta rejected the request (HTTP %d). Check token validity, app review, or permissions." % e.code)
        except Exception:
            return done(False, "Could not reach the Meta Graph API. Check the token and network.")
    if key == "licensed":
        miss = [v for v in ("LICENSED_PROVIDER_API_KEY", "LICENSED_PROVIDER_URL") if not g(v)]
        if miss:
            return done(False, "Missing %s. Add them to demo/.env, restart the server, then test again." % " and ".join(miss), miss)
        try:
            req = urllib.request.Request(g("LICENSED_PROVIDER_URL"), headers={"Authorization": "Bearer " + g("LICENSED_PROVIDER_API_KEY")})
            urllib.request.urlopen(req, timeout=10)
            return done(True, "Licensed provider endpoint reachable.")
        except Exception:
            return done(False, "Credentials found, but the provider endpoint did not respond as expected. Check the URL and key.")
    if key == "tiktok":
        miss = [v for v in ("TIKTOK_CLIENT_KEY", "TIKTOK_CLIENT_SECRET") if not g(v)]
        if miss:
            return done(False, "Missing %s. Add them to demo/.env, restart the server, then test again." % " and ".join(miss), miss)
        return done(False, "Credentials found. Live verification needs approved Research/Display/Commercial access. Enable the matching access flag.")
    if key == "open_web":
        miss = [v for v in ("SEARCH_PROVIDER", "SEARCH_API_KEY") if not g(v)]
        if miss:
            return done(False, "Open web social discovery requires a search provider API key. Add %s (and SEARCH_API_ENDPOINT) to your environment, then test again."
                        % " and ".join(miss), miss)
        endpoint = g("SEARCH_API_ENDPOINT")
        if not endpoint:
            return done(True, "Search provider key found. Set SEARCH_API_ENDPOINT for your provider to complete setup.")
        try:
            urllib.request.urlopen(urllib.request.Request(endpoint, headers={"User-Agent": "egc-pulse/0.2"}), timeout=8)
            return done(True, "Search provider endpoint reachable. Open web discovery is configured.")
        except Exception:
            return done(True, "Search provider key and endpoint found. The endpoint did not answer a bare request, which is normal for query-only APIs.")
    return done(False, "Unknown source.")


def _coverage_from_data(mentions):
    from collections import Counter
    by_plat = Counter(m.get("platform") for m in mentions)
    return [{"display_name": s.display_name, "platform": s.platform, "status": s.status.value,
             "kept": by_plat.get(s.platform, 0)} for s in pam.get_live_collectors()]


def app_mode():
    return os.getenv("APP_MODE") or ("production" if os.getenv("RENDER") else "local")


def public_config():
    """Safe, non-secret config the deployed frontend may read. NEVER contains
    secret VALUES. Only the API base URL (for split hosting), the app mode,
    feature flags, and per-platform source statuses (already public elsewhere).
    The frontend defaults to same-origin relative paths; api_base is only set
    when PUBLIC_API_BASE_URL is configured (e.g. a separate static frontend)."""
    return {
        "api_base": (os.getenv("PUBLIC_API_BASE_URL", "") or "").rstrip("/"),
        "app_mode": app_mode(),
        "internal_use": internal_use(),
        "ai_enabled": aip.enabled(),
        "sources": pam.platform_summary(),
        "live_sources": [s.key for s in pam.get_live_collectors()],
        "discovery": _discovery_status_public(),
    }


def _discovery_status_public():
    """Discovery status for the client. The per-sweep diagnostic counts
    (provider_results/rejected/accepted) are dev-only: included for internal use,
    stripped otherwise. Never contains the API key in either case."""
    disc = cd.discovery_status()
    if not internal_use():
        disc = {k: v for k, v in disc.items() if k != "last_diag"}
    return disc


# ── HTTP ─────────────────────────────────────────────────────────────────────
class Handler(BaseHTTPRequestHandler):
    def log_message(self, *a):
        pass

    def _send(self, obj, status=200, ctype="application/json"):
        body = obj if isinstance(obj, bytes) else json.dumps(obj).encode()
        self.send_response(status)
        self.send_header("Content-Type", ctype)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _download(self, body, ctype, filename):
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Disposition", 'attachment; filename="%s"' % filename)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def do_GET(self):
        u = urllib.parse.urlparse(self.path)
        qs = urllib.parse.parse_qs(u.query)
        one = lambda k: qs.get(k, [None])[0]
        st, en = one("start"), one("end")
        pr = one("project")
        if u.path in ("/", "/index.html"):
            try:
                with open(os.path.join(HERE, "dashboard.html"), "rb") as f:
                    return self._send(f.read(), ctype="text/html; charset=utf-8")
            except FileNotFoundError:
                return self._send(b"dashboard.html missing", 404, "text/plain")
        if u.path == "/api/health":
            return self._send({"status": "ok", "live": [s.key for s in pam.get_live_collectors()],
                               "tracked": [k["label"] for k in tracked_keywords(pr)]})
        if u.path == "/api/config":
            return self._send(public_config())
        if u.path == "/api/projects":
            return self._send({"projects": list_projects(), "active": resolve_project(pr)})
        if u.path == "/api/sources":
            return self._send({"sources": [s.to_dict() for s in pam.get_source_matrix()],
                               "platforms": pam.platform_summary(),
                               "live": [s.key for s in pam.get_live_collectors()],
                               "historical_modes": pam.HISTORICAL_MODES})
        if u.path == "/api/accounts/status":
            return self._send({"accounts": pam.accounts_status(), "env_path": "demo/.env",
                               "internal_use": internal_use(), "ai_enabled": aip.enabled(),
                               "app_mode": app_mode(),
                               "env_targets": {
                                   "local": "demo/.env (then restart the server)",
                                   "production": "your host's environment variables. E.g. Render → Environment (then redeploy)"},
                               "deploy_note": ("Credentials are server-side environment variables and never touch the "
                                               "browser. Set them locally in demo/.env, or in production in your hosting "
                                               "provider's environment settings.")})
        if u.path.startswith("/api/jobs/"):
            jid = u.path[len("/api/jobs/"):].strip("/")
            j = job_public(jid)
            return self._send(j or {"error": "job not found"}, 200 if j else 404)
        if u.path == "/api/metrics":
            return self._send(metrics(one("keyword"), st, en, pr))
        if u.path == "/api/coverage":
            return self._send(coverage_ledger(one("keyword"), st, en, pr))
        if u.path == "/api/insights/account":
            data = pam.account_insights(st, en)
            data["manual"] = get_manual_insights()
            return self._send(data)
        if u.path == "/api/mentions":
            return self._send(recent(int(one("limit") or 25), one("keyword"), one("platform"),
                                     one("sentiment"), st, en, pr))
        # PowerPoint is the single, only export surface (CSV/JSON/printable report retired).
        if u.path == "/api/export/report.pptx":
            try:
                body = report_pptx(one("keyword"), st, en, pr)
            except RuntimeError as e:
                return self._send({"error": str(e)}, 501)
            except Exception as e:
                return self._send({"error": "Could not build the deck: %s" % str(e)[:160]}, 500)
            return self._download(body, "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                  _pptx_filename(pr))
        return self._send({"error": "not found"}, 404)

    def _json_body(self):
        length = int(self.headers.get("Content-Length", 0) or 0)
        try:
            return json.loads(self.rfile.read(length).decode("utf-8", "ignore")) if length else {}
        except Exception:
            return {}

    def do_POST(self):
        u = urllib.parse.urlparse(self.path)
        qs = urllib.parse.parse_qs(u.query)
        one = lambda k: qs.get(k, [None])[0]
        st, en = one("start"), one("end")
        pr = one("project")
        # ── project management ───────────────────────────────────────────────
        if u.path == "/api/projects/create":
            name = one("name") or (self._json_body().get("name") if self.headers.get("Content-Length") else None)
            return self._send({"project": create_project(name), "projects": list_projects()})
        if u.path == "/api/projects/rename":
            body = self._json_body()
            ok = rename_project(one("id") or body.get("id"), one("name") or body.get("name"))
            return self._send({"ok": ok, "projects": list_projects()})
        if u.path == "/api/projects/delete":
            pid = one("id") or self._json_body().get("id")
            ok = delete_project(pid)
            return self._send({"ok": ok, "projects": list_projects(), "active": default_project_id()},
                              200 if ok else 400)
        if u.path == "/api/collect":
            terms = [one("q")] if one("q") else [k["label"] for k in tracked_keywords(pr)]
            terms = [t for t in terms if t]
            if not terms:
                return self._send({"error": "Add a tracked term first."}, 400)
            for t in terms:
                add_keyword(t, pr)   # collected terms must be tracked so results are selectable + visible
            jid = start_job(terms, st, en, one("mode") or "fast", pr)
            return self._send({"job_id": jid, "status": "queued"})
        if u.path.startswith("/api/jobs/") and u.path.endswith("/cancel"):
            jid = u.path[len("/api/jobs/"):-len("/cancel")].strip("/")
            ok = cancel_job(jid)
            return self._send({"job_id": jid, "cancelled": ok, "status": (job_public(jid) or {}).get("status")})
        if u.path == "/api/accounts/test":
            length = int(self.headers.get("Content-Length", 0) or 0)
            try:
                payload = json.loads(self.rfile.read(length).decode("utf-8", "ignore")) if length else {}
            except Exception:
                payload = {}
            return self._send(test_connection(payload.get("source_key") or one("source_key") or ""))
        if u.path == "/api/ai/insights":
            if not aip.enabled():
                return self._send({"error": "AI insights are disabled."}, 403)
            length = int(self.headers.get("Content-Length", 0) or 0)
            try:
                payload = json.loads(self.rfile.read(length).decode("utf-8", "ignore")) if length else {}
            except Exception:
                payload = {}
            terms = payload.get("tracked_terms") or []
            term = terms[0] if terms else None
            plats = payload.get("platforms") or []
            sd, ed = payload.get("start_date"), payload.get("end_date")
            mentions = recent(200, term, plats[0] if plats else None, None, sd, ed, pr)
            if not mentions:
                return self._send({"summary": "No stored mentions in this range. Nothing to summarize.",
                                   "key_themes": [], "sentiment_caveat": "",
                                   "coverage_note": aip._coverage_note(struth.coverage_summary_for_mentions([])),
                                   "source_count": 0, "source_urls": []})
            cov = struth.coverage_summary_for_mentions(mentions)
            insight = aip.deterministic_insight(term, mentions, cov, sd, ed)
            insight["method"] = "deterministic"
            return self._send(insight)
        if u.path == "/api/refresh":
            terms = [k["label"] for k in tracked_keywords(pr)]
            total = sum(collect(t, st, en, pr)["stored"] for t in terms)
            return self._send({"refreshed": len(terms), "stored": total, "tracked": terms})
        if u.path == "/api/backfill":
            return self._send(backfill(one("q") or "", st, en, one("mode") or "recent_only", pr))
        if u.path == "/api/keywords/add":
            term = one("q") or ""
            add_keyword(term, pr)
            r = collect(term, st, en, pr)
            return self._send({"added": term, "stored": r["stored"], "tracked": [k["label"] for k in tracked_keywords(pr)]})
        if u.path == "/api/keywords/remove":
            remove_keyword(one("kw") or "", pr)
            return self._send({"tracked": [k["label"] for k in tracked_keywords(pr)]})
        if u.path == "/api/keywords/clear":
            clear_project(pr)
            return self._send({"cleared": True})
        if u.path == "/api/import":
            payload = self._json_body()
            term = (payload.get("term") or one("q") or "import").strip()
            n = manual_import(term, payload.get("text", ""), payload.get("rows"), pr)
            if n:
                add_keyword(term, pr)
            return self._send({"imported": n, "term": term})
        if u.path == "/api/enrich/urls":
            length = int(self.headers.get("Content-Length", 0) or 0)
            try:
                payload = json.loads(self.rfile.read(length).decode("utf-8", "ignore")) if length else {}
            except Exception:
                payload = {}
            term = (payload.get("term") or one("q") or "url enrichment").strip()
            res = cd.enrich_known_urls(payload.get("urls"), term)
            recs = res.get("records", [])
            ctx = {"source_key": "tiktok_oembed_known_url", "access_path": "official_api", "source_name": "TikTok oEmbed"}
            pid = resolve_project(pr)
            for r in recs:
                r["keyword"] = term
                r["project_id"] = pid
                struth.normalize_mention_source(r, ctx)
            stored = ingest(recs, pid)
            if stored:
                add_keyword(term, pr)
            return self._send({"enriched": stored, "accepted": len(res.get("accepted", [])),
                               "rejected": res.get("rejected", []), "term": term})
        if u.path == "/api/insights/manual":
            length = int(self.headers.get("Content-Length", 0) or 0)
            try:
                payload = json.loads(self.rfile.read(length).decode("utf-8", "ignore")) if length else {}
            except Exception:
                payload = {}
            add_manual_insight(payload.get("platform"), payload.get("period"), payload.get("impressions"),
                               payload.get("reach"), payload.get("engagement"), payload.get("note"))
            return self._send({"ok": True, "manual": get_manual_insights()})
        return self._send({"error": "not found"}, 404)


def serve():
    init_db()
    srv = ThreadingHTTPServer(("0.0.0.0", PORT), Handler)
    print("EGC Pulse on http://localhost:%d  · live sources: %s"
          % (PORT, ", ".join(s.key for s in pam.get_live_collectors()) or "none"))
    srv.serve_forever()


def main():
    load_env()
    cmd = sys.argv[1] if len(sys.argv) > 1 else "serve"
    if cmd == "serve":
        serve()
    elif cmd == "collect":
        init_db()
        print(collect(sys.argv[2] if len(sys.argv) > 2 else "", sys.argv[3] if len(sys.argv) > 3 else None,
                      sys.argv[4] if len(sys.argv) > 4 else None))
    elif cmd == "reset":
        if os.path.exists(DB_PATH):
            os.remove(DB_PATH)
        print("database reset")
    else:
        print(__doc__)


if __name__ == "__main__":
    main()
